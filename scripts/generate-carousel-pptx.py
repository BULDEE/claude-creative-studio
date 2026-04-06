"""
generate-carousel-pptx.py

Generates an editable PPTX carousel from brand.json + carousel.json + copy.md.
All text is real PPTX text (not images) — fully editable in PowerPoint / Keynote.

Usage:
    python3 scripts/generate-carousel-pptx.py \\
        --brand  /path/to/brand.json \\
        --carousel /path/to/carousel.json \\
        --copy   /path/to/copy.md \\
        --output /path/to/carousel.pptx
"""

import argparse
import json
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ModuleNotFoundError:
    print(
        "\n[ERROR] python-pptx is not installed.\n"
        "Install it with:\n"
        "    pip3 install python-pptx\n"
        "or, inside a virtual environment:\n"
        "    python3 -m pip install python-pptx\n"
    )
    sys.exit(1)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

PIXELS_TO_EMU = 9525
SLIDE_WIDTH_PX = 1080
SLIDE_HEIGHT_PX = 1350
FONT_FALLBACK_DISPLAY = "Space Grotesk"
FONT_FALLBACK_BODY = "Inter"
FONT_FALLBACK_MONO = "JetBrains Mono"
FORMAT_GRID_LABELS = ["Carousel", "Vid\u00e9o", "Post court", "Newsletter"]
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---------------------------------------------------------------------------
# EMU + color helpers
# ---------------------------------------------------------------------------


def to_emu(pixels: int | float) -> int:
    """Convert pixels to EMU."""
    return int(pixels * PIXELS_TO_EMU)


def hex_to_rgb(hex_string: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    cleaned = hex_string.lstrip("#")
    if len(cleaned) != 6:
        raise ValueError(f"Invalid hex color: {hex_string!r}")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def _faded(base: RGBColor) -> RGBColor:
    """Return a very dark version of base (SIGNAL background decoration)."""
    return RGBColor(max(0, base[0] // 12), max(0, base[1] // 12), max(0, base[2] // 12))


# ---------------------------------------------------------------------------
# brand.json — color extraction
# ---------------------------------------------------------------------------


def _resolve_color_entry(entry, fallback: str) -> str:
    """Resolve a color entry that is either a dict with 'hex' or a bare string."""
    if isinstance(entry, dict):
        return entry.get("hex", fallback)
    if isinstance(entry, str):
        return entry
    return fallback


def _bg_hex(colors: dict) -> str:
    """Resolve background hex from 'background' or 'neutral.bg'."""
    raw = colors.get("background") or colors.get("neutral", {}).get("bg")
    return _resolve_color_entry(raw, "#0A0A0A") if raw else "#0A0A0A"


def _surface_hex(colors: dict) -> str:
    """Resolve surface hex from 'surface' or 'neutral.surface'."""
    raw = colors.get("surface") or colors.get("neutral", {}).get("surface")
    return _resolve_color_entry(raw, "#1A1A1A") if raw else "#1A1A1A"


def _text_hexes(colors: dict) -> tuple[str, str, str]:
    """Return (white_hex, gray_hex, dark_gray_hex) from brand colors."""
    text = colors.get("text", {})
    white = text.get("primary", "#FFFFFF") if isinstance(text, dict) else "#FFFFFF"
    gray = text.get("secondary", "#A0A0A0") if isinstance(text, dict) else "#A0A0A0"
    dark_gray = colors.get("palette", {}).get("neutral", {}).get("600", "#525252")
    return white, gray, dark_gray


def _brand_color_map(colors: dict) -> dict:
    """Build the flat color map from the brand colors block."""
    pri = hex_to_rgb(_resolve_color_entry(colors.get("primary", {}), "#CCFF00"))
    sec = hex_to_rgb(_resolve_color_entry(colors.get("secondary", {}), "#FF2D8A"))
    acc = hex_to_rgb(_resolve_color_entry(colors.get("accent", {}), "#FFE500"))
    white_h, gray_h, dark_h = _text_hexes(colors)
    return {
        "color_bg": hex_to_rgb(_bg_hex(colors)),
        "color_bg_light": RGBColor(0x11, 0x11, 0x11),
        "color_primary": pri, "color_secondary": sec, "color_accent": acc,
        "color_white": hex_to_rgb(white_h), "color_gray": hex_to_rgb(gray_h),
        "color_dark_gray": hex_to_rgb(dark_h),
        "color_border": hex_to_rgb(_surface_hex(colors)),
        "color_faded_primary": _faded(pri),
        "color_faded_accent": _faded(acc),
        "color_faded_secondary": _faded(sec),
    }


# ---------------------------------------------------------------------------
# brand.json — font extraction + top-level parser
# ---------------------------------------------------------------------------


def _font_family(typography: dict, key: str, fallback: str) -> str:
    """Resolve a font family from a typography entry (dict or bare string)."""
    entry = typography.get(key)
    if isinstance(entry, dict):
        return entry.get("family", fallback)
    if isinstance(entry, str):
        return entry
    return fallback


def parse_brand_json(content: str) -> dict:
    """
    Extract colors, fonts, and brand name from brand.json text.

    Returned keys: color_bg, color_bg_light, color_primary, color_secondary,
    color_accent, color_white, color_gray, color_dark_gray, color_border,
    color_faded_primary, color_faded_accent, color_faded_secondary,
    font_display, font_body, font_mono, brand_name.
    """
    raw = json.loads(content)
    typo = raw.get("typography", {})
    return {
        **_brand_color_map(raw.get("colors", {})),
        "font_display": _font_family(typo, "display", FONT_FALLBACK_DISPLAY),
        "font_body": _font_family(typo, "body", FONT_FALLBACK_BODY),
        "font_mono": _font_family(typo, "mono", FONT_FALLBACK_MONO),
        "brand_name": raw.get("name", "Brand"),
    }


# ---------------------------------------------------------------------------
# copy.md parser
# ---------------------------------------------------------------------------


def _flush_section(sections: dict, slide_num: int | None, lines: list[str]) -> None:
    """Persist accumulated lines into sections when a slide is active."""
    if slide_num is not None:
        sections[slide_num] = lines


def _is_slide_header(line: str) -> re.Match | None:
    return re.match(r"^##\s+Slide\s+(\d+)", line, re.IGNORECASE)


def _is_non_slide_header(line: str) -> bool:
    return bool(re.match(r"^##\s+(?!Slide\s+\d+)", line, re.IGNORECASE))


def _process_copy_line(
    raw_line: str,
    sections: dict,
    current: int | None,
    lines: list[str],
) -> tuple[int | None, list[str]]:
    """Process one raw line from copy.md and return updated (current, lines)."""
    slide_match = _is_slide_header(raw_line)
    if slide_match:
        _flush_section(sections, current, lines)
        return int(slide_match.group(1)), []
    if _is_non_slide_header(raw_line):
        _flush_section(sections, current, lines)
        return None, []
    if raw_line.strip() in ("---", ""):
        return current, lines
    if current is not None:
        stripped = raw_line.strip()
        if stripped:
            lines.append(stripped)
    return current, lines


def parse_copy_md(content: str) -> dict[int, list[str]]:
    """
    Split copy.md into per-slide text.

    Returns dict mapping slide_number → list of non-empty lines under
    that slide's '## Slide N' header. Stops at non-slide '##' sections.
    """
    sections: dict[int, list[str]] = {}
    current: int | None = None
    lines: list[str] = []
    for raw_line in content.splitlines():
        current, lines = _process_copy_line(raw_line, sections, current, lines)
    _flush_section(sections, current, lines)
    return sections


# ---------------------------------------------------------------------------
# Markdown helpers
# ---------------------------------------------------------------------------


def strip_md_bold(text: str) -> str:
    """Remove **bold** markers from a string."""
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def clean_lines(copy_lines: list[str]) -> list[str]:
    """Strip markdown bold from every line in a list."""
    return [strip_md_bold(line) for line in copy_lines]


# ---------------------------------------------------------------------------
# Drawing — textframe primitives
# ---------------------------------------------------------------------------


def _init_frame(slide, left, top, width, height, anchor) -> object:
    """Create a textbox with zeroed margins and return its text frame."""
    box = slide.shapes.add_textbox(to_emu(left), to_emu(top), to_emu(width), to_emu(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = None
    frame.vertical_anchor = anchor
    frame.margin_top = frame.margin_bottom = frame.margin_left = frame.margin_right = 0
    return frame


def _style_run(run, text: str, font: str, size, color: RGBColor, bold: bool) -> None:
    """Apply text content and visual style to a run."""
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _set_spacing(run, letter_spacing: int) -> None:
    """Set DrawingML character spacing on a run element."""
    run.font._element.attrib[f"{{{DRAWINGML_NS}}}spc"] = str(int(letter_spacing))


def _style_paragraph(paragraph, align, size, line_height: float) -> None:
    """Apply paragraph-level formatting."""
    paragraph.alignment = align
    paragraph.line_spacing = Pt(size * line_height)
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def insert_text(
    slide, left, top, width, height, text: str, font: str, size,
    color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT,
    line_height: float = 1.2, letter_spacing: int | None = None,
    anchor=MSO_ANCHOR.TOP,
) -> object:
    """Add a single-paragraph text box. All positions in pixels."""
    frame = _init_frame(slide, left, top, width, height, anchor)
    para = frame.paragraphs[0]
    _style_paragraph(para, align, size, line_height)
    run = para.add_run()
    _style_run(run, text, font, size, color, bold)
    if letter_spacing is not None:
        _set_spacing(run, letter_spacing)
    return frame


# ---------------------------------------------------------------------------
# Drawing — multiline text
# ---------------------------------------------------------------------------


def _append_paragraph(frame, index: int, text: str, color: RGBColor, bold: bool,
                       font: str, size, align, line_height: float, spacing) -> None:
    """Append one styled paragraph to a text frame."""
    para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
    _style_paragraph(para, align, size, line_height)
    para.space_after = spacing
    run = para.add_run()
    _style_run(run, text, font, size, color, bold)


def insert_multiline(
    slide, left, top, width, height,
    lines: list[tuple[str, RGBColor, bool]],
    font: str, size, align=PP_ALIGN.LEFT,
    line_height: float = 1.2, para_spacing=None,
) -> object:
    """Add a multi-paragraph text box. lines = list of (text, color, bold)."""
    frame = _init_frame(slide, left, top, width, height, MSO_ANCHOR.TOP)
    spacing = para_spacing if para_spacing is not None else Pt(size * 0.5)
    for idx, (txt, color, bold) in enumerate(lines):
        _append_paragraph(frame, idx, txt, color, bold, font, size, align, line_height, spacing)
    return frame


# ---------------------------------------------------------------------------
# Drawing — shapes
# ---------------------------------------------------------------------------


def _fill_shape(shape, color: RGBColor | None) -> None:
    """Apply solid fill or transparent fill to a shape."""
    if color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()


def _border_shape(shape, color: RGBColor | None) -> None:
    """Apply border or transparent border to a shape."""
    if color is not None:
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def _add_shape(slide, shape_type, left, top, width, height):
    """Add an MSO shape and return it."""
    return slide.shapes.add_shape(
        shape_type, to_emu(left), to_emu(top), to_emu(width), to_emu(height),
    )


def insert_rectangle(slide, left, top, width, height,
                     fill_color: RGBColor, border_color: RGBColor | None = None) -> object:
    """Add a plain filled rectangle."""
    shape = _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_shape(shape, fill_color)
    _border_shape(shape, border_color)
    return shape


def insert_rounded_rect(slide, left, top, width, height,
                        fill_color: RGBColor | None = None,
                        border_color: RGBColor | None = None) -> object:
    """Add a rounded rectangle with optional fill and border."""
    shape = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_shape(shape, fill_color)
    _border_shape(shape, border_color)
    return shape


# ---------------------------------------------------------------------------
# Background
# ---------------------------------------------------------------------------


def apply_background(slide, color: RGBColor) -> None:
    """Flood-fill a slide with a solid background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Chrome — sub-elements
# ---------------------------------------------------------------------------


def _chrome_counter(slide, slide_num: int, total: int, brand: dict) -> None:
    """Render top-right slide counter."""
    insert_text(
        slide, 800, 50, 200, 24, f"{slide_num:02d} / {total:02d}",
        brand["font_mono"], 14, brand["color_dark_gray"],
        align=PP_ALIGN.RIGHT, letter_spacing=100,
    )


def _chrome_watermark(slide, brand: dict, size: int = 16) -> None:
    """Render bottom-right brand watermark."""
    insert_text(
        slide, 800, 1270, 200, 24, brand["brand_name"],
        brand["font_display"], size, brand["color_primary"],
        bold=True, align=PP_ALIGN.RIGHT,
    )


def _chrome_swipe(slide, brand: dict) -> None:
    """Render bottom-left 'SWIPE →' label."""
    insert_text(
        slide, 80, 1270, 200, 24, "SWIPE \u2192",
        brand["font_mono"], 12, brand["color_dark_gray"], letter_spacing=100,
    )


def add_chrome(slide, slide_num: int, total: int, brand: dict, show_swipe: bool = True) -> None:
    """Render slide counter, watermark, and optional swipe indicator."""
    _chrome_counter(slide, slide_num, total, brand)
    _chrome_watermark(slide, brand)
    if show_swipe:
        _chrome_swipe(slide, brand)


# ---------------------------------------------------------------------------
# HOOK slide
# ---------------------------------------------------------------------------


def _hook_texts(copy_lines: list[str]) -> tuple[str, str]:
    """Return (title, subtitle) from HOOK copy."""
    title_lines = [strip_md_bold(ln) for ln in copy_lines if ln.startswith("**")]
    sub_lines = [strip_md_bold(ln) for ln in copy_lines if not ln.startswith("**")]
    title = title_lines[0] if title_lines else (copy_lines[0] if copy_lines else "")
    return title, (sub_lines[0] if sub_lines else "")


def build_hook_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """HOOK: brand tag + large title + subtitle."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    insert_text(
        slide, 80, 420, 500, 20,
        f"{brand['brand_name'].upper()} \u2014 CAROUSEL",
        brand["font_mono"], 13, brand["color_primary"], letter_spacing=150,
    )
    title, subtitle = _hook_texts(copy_lines)
    insert_text(
        slide, 80, 470, 800, 280, title,
        brand["font_display"], 72, brand["color_white"], bold=True, line_height=1.08,
    )
    if subtitle:
        insert_text(
            slide, 80, 760, 700, 80, subtitle,
            brand["font_body"], 28, brand["color_gray"], line_height=1.5,
        )


# ---------------------------------------------------------------------------
# DATA slide
# ---------------------------------------------------------------------------


def _data_parse_numbers(copy_lines: list[str]) -> tuple[list, list, list]:
    """Return (number_tokens, label_texts, punchline_texts) from DATA copy."""
    nums: list[str] = []
    labels: list[str] = []
    punches: list[str] = []
    for line in copy_lines:
        cleaned = strip_md_bold(line)
        match = re.match(r"^(\d+\s*[%]?)(.*)$", cleaned)
        if match:
            nums.append(match.group(1).replace(" ", ""))
            if match.group(2).strip():
                labels.append(match.group(2).strip())
        elif cleaned:
            punches.append(cleaned)
    if not nums:
        for line in copy_lines:
            bold_match = re.search(r"\*\*(\d+\s*[%]?)\*\*", line)
            if bold_match:
                nums.append(bold_match.group(1).replace(" ", ""))
                labels.append(re.sub(r"\*\*\d+\s*[%]?\*\*\s*", "", line).strip())
    return nums, labels, punches


def _data_stat(slide, brand, num: str, label: str, color: RGBColor, top_n: int, top_l: int) -> None:
    """Render one big number + label row."""
    insert_text(
        slide, 80, top_n, 920, 220, num,
        brand["font_mono"], 180, color, bold=True, align=PP_ALIGN.CENTER, line_height=1.0,
    )
    if label:
        insert_text(
            slide, 80, top_l, 920, 40, label,
            brand["font_body"], 24, brand["color_gray"], align=PP_ALIGN.CENTER,
        )


def build_data_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """DATA: two big stats + divider + punchline."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    nums, labels, punches = _data_parse_numbers(copy_lines)
    _data_stat(slide, brand, nums[0] if nums else "?", labels[0] if labels else "",
               brand["color_primary"], 140, 370)
    insert_rectangle(slide, 480, 450, 120, 2, brand["color_dark_gray"])
    _data_stat(slide, brand, nums[1] if len(nums) > 1 else "?", labels[1] if len(labels) > 1 else "",
               brand["color_secondary"], 500, 730)
    if punches:
        insert_text(
            slide, 80, 830, 920, 40, punches[0],
            brand["font_body"], 22, brand["color_white"], bold=True, align=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# PAIN slide
# ---------------------------------------------------------------------------


def _pain_tuples(body: list[str], brand: dict) -> list[tuple[str, RGBColor, bool]]:
    """Color-code PAIN lines: lines with 'non' get secondary color."""
    return [
        (ln, brand["color_secondary"] if re.search(r"\bnon\b", ln, re.IGNORECASE)
         else brand["color_white"], False)
        for ln in body
    ]


def build_pain_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """PAIN: left accent bar + color-coded lines + footer."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    insert_rectangle(slide, 76, 400, 4, 380, brand["color_secondary"])
    stripped = clean_lines(copy_lines)
    footer = stripped[-1] if len(stripped) > 1 else ""
    body = stripped[:-1] if len(stripped) > 1 else stripped
    tuples = _pain_tuples(body, brand)
    if tuples:
        insert_multiline(
            slide, 120, 400, 800, 380, tuples,
            brand["font_body"], 36, line_height=1.3, para_spacing=Pt(20),
        )
    if footer:
        insert_text(slide, 120, 840, 700, 30, footer, brand["font_body"], 20, brand["color_dark_gray"])


# ---------------------------------------------------------------------------
# SHIFT slide
# ---------------------------------------------------------------------------


def _shift_badge_word(copy_lines: list[str]) -> str:
    """Extract the first bold token from copy as the badge word."""
    for raw in copy_lines:
        match = re.search(r"\*\*(.+?)\*\*", raw)
        if match:
            return match.group(1).upper()
    return "CONCEPT"


def _shift_render_badge(slide, brand: dict, word: str) -> None:
    """Render the colored keyword badge."""
    insert_rectangle(slide, 80, 620, 400, 90, brand["color_primary"])
    insert_text(
        slide, 80, 628, 400, 80, word,
        brand["font_display"], 64, brand["color_bg"],
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def build_shift_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """SHIFT: problem → bridge → badge → resolution."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    stripped = clean_lines(copy_lines)
    problem = stripped[0] if stripped else ""
    bridge = " ".join(stripped[1:-1]) if len(stripped) > 2 else ""
    resolution = stripped[-1] if len(stripped) > 1 else ""
    if problem:
        insert_text(slide, 80, 380, 700, 50, problem, brand["font_body"], 28, brand["color_gray"], line_height=1.4)
    if bridge:
        insert_text(slide, 80, 450, 700, 120, bridge, brand["font_body"], 28, brand["color_white"], line_height=1.5)
    _shift_render_badge(slide, brand, _shift_badge_word(copy_lines))
    if resolution:
        insert_text(slide, 80, 750, 600, 40, resolution, brand["font_body"], 26, brand["color_primary"], bold=True)


# ---------------------------------------------------------------------------
# SIGNAL shared helpers
# ---------------------------------------------------------------------------


def _signal_colors(signal_number: int, brand: dict) -> tuple[RGBColor, RGBColor]:
    """Return (label_color, faded_color) for a signal number."""
    if signal_number <= 2:
        return brand["color_primary"], brand["color_faded_primary"]
    if signal_number == 3:
        return brand["color_accent"], brand["color_faded_accent"]
    return brand["color_secondary"], brand["color_faded_secondary"]


def _signal_header(slide, brand: dict, num: int, label_color: RGBColor, faded: RGBColor) -> None:
    """Render the faded background number and SIGNAL label."""
    insert_text(
        slide, 70, 70, 300, 140, f"{num:02d}",
        brand["font_mono"], 120, faded, bold=True, line_height=1.0,
    )
    insert_text(slide, 80, 350, 200, 20, "SIGNAL", brand["font_mono"], 14, label_color, letter_spacing=150)


def _signal_title(slide, brand: dict, title: str) -> None:
    """Render the SIGNAL slide main title."""
    insert_text(
        slide, 80, 390, 750, 180, title,
        brand["font_display"], 48, brand["color_white"], bold=True, line_height=1.15,
    )


def _signal_bullets(bullet_lines: list[str], brand: dict) -> list[tuple[str, RGBColor, bool]]:
    """Build arrow-prefixed bullet tuples; last bullet white+bold."""
    last = len(bullet_lines) - 1
    return [
        (f"\u2192  {ln}", brand["color_white"] if idx == last else brand["color_gray"], idx == last)
        for idx, ln in enumerate(bullet_lines)
    ]


# ---------------------------------------------------------------------------
# SIGNAL default slide
# ---------------------------------------------------------------------------


def build_signal_slide(
    slide, brand: dict, copy_lines: list[str], slide_num: int, total: int, signal_number: int,
) -> None:
    """SIGNAL N (default): faded number + label + title + arrow bullets."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    label_color, faded = _signal_colors(signal_number, brand)
    _signal_header(slide, brand, signal_number, label_color, faded)
    stripped = clean_lines(copy_lines)
    _signal_title(slide, brand, stripped[0] if stripped else "")
    bullets = _signal_bullets(stripped[1:], brand)
    if bullets:
        insert_multiline(
            slide, 80, 610, 750, 200, bullets,
            brand["font_body"], 22, line_height=1.4, para_spacing=Pt(12),
        )


# ---------------------------------------------------------------------------
# SIGNAL 02 — format grid
# ---------------------------------------------------------------------------


def _signal02_grid(slide, brand: dict) -> None:
    """Render the 4-card 2x2 format grid."""
    col1, col2, card_w, card_h, gap = 80, 370, 260, 80, 20
    top = 630
    colors = [brand["color_primary"], brand["color_secondary"], brand["color_accent"], brand["color_white"]]
    for idx, (label, color) in enumerate(zip(FORMAT_GRID_LABELS, colors)):
        card_x = col1 if idx % 2 == 0 else col2
        card_y = top + (idx // 2) * (card_h + gap)
        insert_rounded_rect(slide, card_x, card_y, card_w, card_h, border_color=brand["color_border"])
        insert_text(
            slide, card_x, card_y, card_w, card_h, label,
            brand["font_mono"], 14, color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )


def build_signal_02_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """SIGNAL 02: format grid with 4 rounded-rect cards."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    label_color, faded = _signal_colors(2, brand)
    _signal_header(slide, brand, 2, label_color, faded)
    stripped = clean_lines(copy_lines)
    _signal_title(slide, brand, stripped[0] if stripped else "")
    _signal02_grid(slide, brand)
    if len(stripped) > 1:
        insert_text(slide, 80, 870, 700, 40, stripped[-1], brand["font_body"], 22, brand["color_gray"])


# ---------------------------------------------------------------------------
# SIGNAL 03 — hook comparison rows
# ---------------------------------------------------------------------------


def _signal03_rows(brand: dict) -> list[tuple]:
    """Return illustrative hook comparison row data."""
    return [
        ("Hook variante A", "6.2%", brand["color_gray"], brand["color_border"], False),
        ("Hook variante B", "9.1% \u2605", brand["color_primary"], brand["color_primary"], True),
        ("Hook variante C", "5.8%", brand["color_gray"], brand["color_border"], False),
    ]


def _signal03_render_rows(slide, brand: dict, top: int) -> None:
    """Render comparison rows for SIGNAL 03."""
    row_h, gap = 60, 12
    for idx, (label, score, color, border, winner) in enumerate(_signal03_rows(brand)):
        row_y = top + idx * (row_h + gap)
        insert_rounded_rect(slide, 80, row_y, 700, row_h, border_color=border)
        insert_text(
            slide, 100, row_y, 400, row_h, label,
            brand["font_body"], 18, color, bold=winner, anchor=MSO_ANCHOR.MIDDLE,
        )
        insert_text(
            slide, 550, row_y, 210, row_h, score,
            brand["font_mono"], 14, color, bold=winner, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )


def build_signal_03_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """SIGNAL 03: hook comparison table with scored rows."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    label_color, faded = _signal_colors(3, brand)
    _signal_header(slide, brand, 3, label_color, faded)
    stripped = clean_lines(copy_lines)
    _signal_title(slide, brand, stripped[0] if stripped else "")
    _signal03_render_rows(slide, brand, top=620)
    if len(stripped) > 1:
        insert_text(
            slide, 80, 860, 700, 40, stripped[-1],
            brand["font_body"], 22, brand["color_white"], bold=True,
        )


# ---------------------------------------------------------------------------
# ACTION slide
# ---------------------------------------------------------------------------


def _action_parse(stripped: list[str]) -> tuple[str, list[str], str]:
    """Split ACTION copy into (title, steps, footer)."""
    title = stripped[0] if stripped else ""
    steps: list[str] = []
    footer = ""
    for idx, line in enumerate(stripped[1:], start=1):
        if idx == len(stripped) - 1 and not steps:
            footer = line
        else:
            steps.append(line)
    if not footer and steps:
        footer = steps.pop()
    return title, steps, footer


def build_action_slide(slide, brand: dict, copy_lines: list[str], slide_num: int, total: int) -> None:
    """ACTION: label + title + steps + footer."""
    apply_background(slide, brand["color_bg_light"])
    add_chrome(slide, slide_num, total, brand)
    insert_text(
        slide, 80, 360, 920, 20, "CETTE SEMAINE",
        brand["font_mono"], 14, brand["color_primary"], align=PP_ALIGN.CENTER, letter_spacing=150,
    )
    title, steps, footer = _action_parse(clean_lines(copy_lines))
    insert_text(
        slide, 80, 410, 920, 70, title,
        brand["font_display"], 44, brand["color_white"], bold=True, align=PP_ALIGN.CENTER, line_height=1.1,
    )
    if steps:
        insert_multiline(
            slide, 160, 550, 760, 280, [(ln, brand["color_white"], False) for ln in steps],
            brand["font_body"], 24, align=PP_ALIGN.LEFT, line_height=1.4, para_spacing=Pt(16),
        )
    if footer:
        insert_text(slide, 80, 890, 920, 40, footer, brand["font_body"], 22, brand["color_dark_gray"],
                    align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# CTA slide
# ---------------------------------------------------------------------------


def _cta_parse(copy_lines: list[str], url: str) -> tuple[str, list[str]]:
    """Return (title, feature_lines) from CTA copy, skipping the URL line."""
    title = ""
    features: list[str] = []
    for line in copy_lines:
        cleaned = strip_md_bold(line)
        if url and cleaned.lower() == url.lower():
            continue
        if not title and line.startswith("**"):
            title = cleaned
        elif cleaned:
            features.append(cleaned)
    return title, features


def _cta_chrome(slide, slide_num: int, total: int, brand: dict) -> None:
    """Render CTA chrome: counter + larger watermark (no swipe)."""
    _chrome_counter(slide, slide_num, total, brand)
    _chrome_watermark(slide, brand, size=24)


def build_cta_slide(
    slide, brand: dict, copy_lines: list[str], slide_num: int, total: int,
    cta_url: str, cta_tagline: str,
) -> None:
    """CTA: title + features + big URL + tagline + larger watermark."""
    apply_background(slide, brand["color_bg"])
    _cta_chrome(slide, slide_num, total, brand)
    title, features = _cta_parse(copy_lines, cta_url)
    insert_text(
        slide, 80, 340, 920, 60, title,
        brand["font_display"], 36, brand["color_white"], bold=True, align=PP_ALIGN.CENTER, line_height=1.2,
    )
    if features:
        insert_multiline(
            slide, 80, 460, 920, 220, [(ln, brand["color_gray"], False) for ln in features],
            brand["font_body"], 22, align=PP_ALIGN.CENTER, line_height=1.4, para_spacing=Pt(10),
        )
    insert_text(
        slide, 80, 760, 920, 80, cta_url,
        brand["font_mono"], 56, brand["color_primary"], bold=True, align=PP_ALIGN.CENTER, line_height=1.0,
    )
    if cta_tagline:
        insert_text(slide, 80, 870, 920, 40, cta_tagline, brand["font_body"], 20, brand["color_dark_gray"],
                    align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide dispatcher
# ---------------------------------------------------------------------------


def _signal_number(role: str) -> int | None:
    """Extract signal number from role string like 'SIGNAL 01'."""
    match = re.search(r"SIGNAL\s*(\d+)", role, re.IGNORECASE)
    return int(match.group(1)) if match else None


def _fallback_slide(slide, brand: dict, role: str, copy_lines: list[str], slide_num: int, total: int) -> None:
    """Render an unrecognised role with a generic layout."""
    apply_background(slide, brand["color_bg"])
    add_chrome(slide, slide_num, total, brand)
    insert_text(
        slide, 80, 400, 920, 60, role,
        brand["font_display"], 36, brand["color_primary"], bold=True, align=PP_ALIGN.CENTER,
    )
    for idx, copy_line in enumerate(copy_lines[:4]):
        insert_text(
            slide, 80, 500 + idx * 60, 920, 50, strip_md_bold(copy_line),
            brand["font_body"], 24, brand["color_white"], align=PP_ALIGN.CENTER,
        )


def _dispatch_signal(slide, brand, copy_lines, slide_num, total, sig_num: int) -> None:
    """Route to the correct SIGNAL variant builder."""
    if sig_num == 2:
        build_signal_02_slide(slide, brand, copy_lines, slide_num, total)
    elif sig_num == 3:
        build_signal_03_slide(slide, brand, copy_lines, slide_num, total)
    else:
        build_signal_slide(slide, brand, copy_lines, slide_num, total, sig_num)


def build_slide(
    deck, layout, slide_entry: dict, brand: dict,
    copy_sections: dict, total_slides: int, cta_url: str, cta_tagline: str,
) -> None:
    """Add one slide to the deck by dispatching on role."""
    slide = deck.slides.add_slide(layout)
    slide_num: int = slide_entry["slide"]
    role: str = slide_entry.get("role", "").upper()
    copy_lines = copy_sections.get(slide_num, [])
    sig = _signal_number(role)
    dispatch = {
        "HOOK": lambda: build_hook_slide(slide, brand, copy_lines, slide_num, total_slides),
        "DATA": lambda: build_data_slide(slide, brand, copy_lines, slide_num, total_slides),
        "PAIN": lambda: build_pain_slide(slide, brand, copy_lines, slide_num, total_slides),
        "SHIFT": lambda: build_shift_slide(slide, brand, copy_lines, slide_num, total_slides),
        "ACTION": lambda: build_action_slide(slide, brand, copy_lines, slide_num, total_slides),
        "CTA": lambda: build_cta_slide(slide, brand, copy_lines, slide_num, total_slides, cta_url, cta_tagline),
    }
    if sig is not None:
        _dispatch_signal(slide, brand, copy_lines, slide_num, total_slides, sig)
    elif role in dispatch:
        dispatch[role]()
    else:
        _fallback_slide(slide, brand, role, copy_lines, slide_num, total_slides)


# ---------------------------------------------------------------------------
# File I/O
# ---------------------------------------------------------------------------


def read_file(path: str) -> str:
    """Read a UTF-8 file or exit with a clear error message."""
    try:
        with open(path, encoding="utf-8") as fh:
            return fh.read()
    except FileNotFoundError:
        print(f"[ERROR] File not found: {path}")
        sys.exit(1)
    except OSError as err:
        print(f"[ERROR] Could not read {path}: {err}")
        sys.exit(1)


# ---------------------------------------------------------------------------
# CLI and orchestration
# ---------------------------------------------------------------------------


def _cli_parser() -> argparse.ArgumentParser:
    """Return the configured argument parser."""
    parser = argparse.ArgumentParser(
        prog="generate-carousel-pptx",
        description="Generate an editable PPTX carousel from brand.json, carousel.json, and copy.md.",
    )
    parser.add_argument("--brand", required=True, metavar="PATH", help="Path to brand.json")
    parser.add_argument("--carousel", required=True, metavar="PATH", help="Path to carousel.json")
    parser.add_argument("--copy", required=True, metavar="PATH", help="Path to copy.md")
    parser.add_argument("--output", required=True, metavar="PATH", help="Destination .pptx path")
    return parser


def _log_brand(brand: dict, path: str) -> None:
    """Print brand summary to stdout."""
    print(f"Reading brand       : {path}")
    print(f"  Brand name        : {brand['brand_name']}")
    print(f"  Display font      : {brand['font_display']}")
    print(f"  Body font         : {brand['font_body']}")
    print(f"  Mono font         : {brand['font_mono']}")


def _log_carousel(carousel: dict, path: str) -> None:
    """Print carousel summary to stdout."""
    print(f"Reading carousel    : {path}")
    print(f"  Title             : {carousel.get('title', '')}")
    print(f"  Slides            : {len(carousel.get('structure', []))}")
    print(f"  CTA URL           : {carousel.get('cta', {}).get('url', '')}")


def _load_inputs(args) -> tuple:
    """Parse all three input files and return (brand, structure, cta_url, cta_tagline, copy_sections)."""
    brand = parse_brand_json(read_file(args.brand))
    _log_brand(brand, args.brand)
    carousel = json.loads(read_file(args.carousel))
    _log_carousel(carousel, args.carousel)
    cta = carousel.get("cta", {})
    copy_sections = parse_copy_md(read_file(args.copy))
    print(f"Reading copy        : {args.copy}")
    print(f"  Parsed sections   : {sorted(copy_sections.keys())}")
    return brand, carousel.get("structure", []), cta.get("url", ""), cta.get("tagline", ""), copy_sections


def _assemble_deck(brand: dict, structure: list, cta_url: str, cta_tagline: str, copy_sections: dict) -> object:
    """Build and return a populated Presentation object."""
    deck = Presentation()
    deck.slide_width = to_emu(SLIDE_WIDTH_PX)
    deck.slide_height = to_emu(SLIDE_HEIGHT_PX)
    layout = deck.slide_layouts[6]
    print(f"\nBuilding {len(structure)} slides...")
    for entry in structure:
        print(f"  Slide {entry.get('slide', '?'):>2}  [{entry.get('role', 'UNKNOWN')}]")
        build_slide(deck, layout, entry, brand, copy_sections, len(structure), cta_url, cta_tagline)
    return deck


def _save_deck(deck, output_path: str) -> None:
    """Save the deck to disk and print a summary."""
    try:
        deck.save(output_path)
    except OSError as err:
        print(f"\n[ERROR] Could not save PPTX to {output_path}: {err}")
        sys.exit(1)
    print(f"\nPPTX saved          : {output_path}")
    print(f"Slides              : {len(deck.slides)}")
    print(f"Dimensions          : {SLIDE_WIDTH_PX}x{SLIDE_HEIGHT_PX}px")
    print("\nNote: Install Space Grotesk, Inter, JetBrains Mono for correct font rendering.")


def main() -> None:
    """Entry point: parse args, load inputs, build, save."""
    args = _cli_parser().parse_args()
    brand, structure, cta_url, cta_tagline, copy_sections = _load_inputs(args)
    deck = _assemble_deck(brand, structure, cta_url, cta_tagline, copy_sections)
    _save_deck(deck, args.output)


if __name__ == "__main__":
    main()
